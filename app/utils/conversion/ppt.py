from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

async def generate_presentation(lesson_plan_data, output_ppt_path):
    """
    Converts lesson plan data into a PowerPoint presentation with appropriate slides for each key-value pair.
    
    :param lesson_plan_data: The lesson plan data to be added to the PowerPoint slides.
    :param output_ppt_path: The file path to save the generated PowerPoint presentation.
    """
    try:
        # Step 1: Create a PowerPoint presentation object
        presentation = Presentation()
        slide_count = 0

        def add_slide_for_content(title, content):
            nonlocal slide_count
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
            slide.shapes.title.text = title.replace("_", " ")  # Remove underscores and replace with spaces

            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content

            # Reduce font size for the content and apply formatting
            for paragraph in content_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    run.font.bold = True  # Example of bold text
                    run.font.italic = True  # Example of italic text

            # Adjust text alignment
            for paragraph in content_placeholder.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT  # You can change alignment here

            slide_count += 1

        # Process each lesson plan
        for lesson in lesson_plan_data["lesson_plan"]:
            # Add main lesson topic slide
            add_slide_for_content("Lesson Topic", lesson.get("Lesson_Topic", ""))

            # Add remaining keys as slides
            for key, value in lesson.items():
                if key != "Lesson_Topic":
                    if isinstance(value, dict):  # For nested dictionaries
                        for sub_key, sub_value in value.items():
                            add_slide_for_content(sub_key.replace("_", " "), sub_value)
                    elif isinstance(value, list):  # For lists
                        add_slide_for_content(key.replace("_", " "), "\n".join(value))
                    else:  # For plain text values
                        add_slide_for_content(key.replace("_", " "), value)

        # Save the presentation to a file
        presentation.save(output_ppt_path)
        print(f"Presentation saved as {output_ppt_path}")

    except Exception as e:
        print(f"An error occurred while generating the presentation: {e}")
        raise


async def refine_ppt_content(output_ppt_path):
    """
    Refines the PowerPoint content by adjusting fonts, sizes, etc.
    
    :param output_ppt_path: Path to the PowerPoint file to be refined.
    """
    try:
        presentation = Presentation(output_ppt_path)

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(16)  # Example: increase font size
                            run.font.bold = True    # Apply bold to all runs
                            run.font.italic = False  # Remove italic formatting

        presentation.save(output_ppt_path)
        print(f"Refined PowerPoint saved as {output_ppt_path}")

    except Exception as e:
        print(f"An error occurred while refining the PowerPoint: {e}")
        raise
